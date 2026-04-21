"""
NVIDIA NIM service for:
- Document OCR & Parsing (NeMo Retriever Parse)
- Text Embeddings (nv-embedqa-e5-v5)
- Chat Completion / RAG Generation (Llama via NIM)
"""

import base64
import io
import json
import math
import re
from typing import Any, Dict, List, Optional

import httpx

from src.core.config import get_settings
from src.core.logger import logger

# Optional dependency
# from openai import AsyncOpenAI  <- Moved to lazy import inside function


settings = get_settings()


def _get_nim_client():
    """Return an AsyncOpenAI client pointed at NVIDIA NIM."""
    try:
        from openai import AsyncOpenAI

        return AsyncOpenAI(
            api_key=settings.NVIDIA_API_KEY,
            base_url=settings.NIM_BASE_URL,
        )
    except ImportError:
        logger.error("NVIDIA AI features require 'openai' library. Run 'uv add openai'.")
        raise


# ─────────────────────────────────────────────
# Document Parsing (NeMo Retriever Parse NIM)
# ─────────────────────────────────────────────


async def _convert_pdf_to_images(file_bytes: bytes) -> List[tuple[int, str]]:
    """
    Converts a PDF (bytes) into a list of (page_num, base64_png) tuples.
    Uses PyMuPDF for high-quality rendering.
    """
    import fitz  # Lazy import to avoid environment crashes

    images = []
    # Open the PDF from bytes
    doc = fitz.open(stream=file_bytes, filetype="pdf")

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        # Render page to a pixmap (DPI 200 for high OCR quality and reliability)
        pix = page.get_pixmap(matrix=fitz.Matrix(200 / 72, 200 / 72))

        # Convert pixmap to PNG bytes
        img_bytes = pix.tobytes("png")
        b64_img = base64.b64encode(img_bytes).decode("utf-8")
        images.append((page_num + 1, b64_img))

    doc.close()
    return images


async def _call_parse_nim(b64_image: str, page_num: int, filename: str) -> List[dict]:
    """
    Calls the NeMo Retriever Parse NIM for a single base64-encoded image.
    """
    payload = {
        "messages": [
            {
                "role": "user",
                "content": [
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:image/png;base64,{b64_image}"},
                    },
                ],
            }
        ],
        "tools": [{"type": "function", "function": {"name": "markdown_no_bbox"}}],
        "model": settings.NIM_PARSE_MODEL,
        "max_tokens": 3000,
        "temperature": 0.0,
    }

    client = _get_nim_client()

    try:
        response = await client.chat.completions.create(
            model=settings.NIM_PARSE_MODEL,
            messages=payload["messages"],
            tools=payload["tools"],
            max_tokens=3000,
            temperature=0.0,
        )

        message = response.choices[0].message
        raw_text = message.content

        # NeMo Retriever Parse often returns content in tool_calls arguments if tools were provided
        if not raw_text and message.tool_calls:
            logger.info(
                f"NIM Parse returned content in tool_calls for '{filename}' page {page_num}."
            )
            raw_text = message.tool_calls[0].function.arguments

        # Log raw response for debugging if it's still empty or suspicious
        if not raw_text:
            logger.warning(
                f"NIM Parse returned empty content for '{filename}' page {page_num}. Response: {response}"
            )
        else:
            logger.info(f"NIM Parse raw output (first 200 chars): {raw_text[:200]}")

    except Exception as e:
        logger.error(f"NIM Parse Exception: {e}")
        raise

    if not raw_text:
        return []

    # 1. Try to find a JSON array [...]
    json_match = re.search(r"\[\s*{.*}\s*\]", raw_text, re.DOTALL)
    if json_match:
        try:
            elements = json.loads(json_match.group())
            for el in elements:
                el["page"] = page_num
            return elements
        except Exception:
            pass

    # 2. Try to find a JSON object { ... "elements": [...] ... }
    obj_match = re.search(r"\{.*\}", raw_text, re.DOTALL)
    if obj_match:
        try:
            data = json.loads(obj_match.group())
            # Check for common result keys
            elements = data.get("elements") or data.get("pages", [{}])[0].get("elements")
            if elements and isinstance(elements, list):
                for el in elements:
                    el["page"] = page_num
                return elements
        except Exception:
            pass

    # 3. Fallback: Treat as raw Markdown or Text
    logger.info(
        f"NIM Parse returned raw text for '{filename}' page {page_num}. Converting to elements."
    )
    # Split by double newline to create pseudo-elements
    paragraphs = [p.strip() for p in raw_text.split("\n\n") if p.strip()]
    return [{"type": "paragraph", "text": p, "page": page_num} for p in paragraphs]


async def parse_document_bytes(file_bytes: bytes, filename: str) -> List[dict]:
    """
    Universal document parsing router. Routes files to specialized parsers
    based on extension.
    """
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "bin"

    if ext == "pdf":
        return await _parse_pdf(file_bytes, filename)
    elif ext in ["png", "jpg", "jpeg", "tiff", "bmp"]:
        return await _parse_visual_image(file_bytes, filename)
    elif ext == "docx":
        return _parse_docx(file_bytes, filename)
    elif ext in ["xlsx", "xls"]:
        return _parse_xlsx(file_bytes, filename)
    elif ext == "pptx":
        return _parse_pptx(file_bytes, filename)
    else:
        # Fallback for text-based formats (TXT, MD, CSV, JSON, Log, etc.)
        return _parse_text_generic(file_bytes, filename, ext)


async def _parse_pdf(file_bytes: bytes, filename: str) -> List[dict]:
    """Renders PDF to images and uses NIM Parse for layout-aware extraction."""
    logger.info(f"Rendering PDF '{filename}' to images for visual parsing...")
    images = await _convert_pdf_to_images(file_bytes)
    logger.info(f"Processing {len(images)} pages for '{filename}'...")

    all_elements = []
    for page_num, b64_img in images:
        logger.info(f"Parsing page {page_num}/{len(images)} of '{filename}'...")
        elements = await _call_parse_nim(b64_img, page_num, filename)
        all_elements.extend(elements)
    return all_elements


async def _parse_visual_image(file_bytes: bytes, filename: str) -> List[dict]:
    """Uses NIM Parse for direct image OCR/Layout analysis."""
    logger.info(f"Processing image document '{filename}'...")
    b64_content = base64.b64encode(file_bytes).decode("utf-8")
    return await _call_parse_nim(b64_content, 1, filename)


def _parse_docx(file_bytes: bytes, filename: str) -> List[dict]:
    """Extracts text, headings, and tables from Word documents using python-docx."""
    from docx import Document

    logger.info(f"Extracting structured text from Word document '{filename}'...")

    doc = Document(io.BytesIO(file_bytes))
    elements = []

    # Process paragraphs with style awareness
    for i, para in enumerate(doc.paragraphs):
        text = para.text.strip()
        if not text:
            continue

        el_type = "paragraph"
        if para.style.name.startswith("Heading"):
            el_type = "heading"

        elements.append({"type": el_type, "text": text, "page": 1})

    # Process tables — convert to simple markdown-like text
    for i, table in enumerate(doc.tables):
        table_data = []
        for row in table.rows:
            row_text = " | ".join([cell.text.strip() for cell in row.cells])
            table_data.append(row_text)

        if table_data:
            elements.append({"type": "table", "text": "\n".join(table_data), "page": 1})

    return elements


def _parse_xlsx(file_bytes: bytes, filename: str) -> List[dict]:
    """Extracts Excel data and converts sheets to Markdown tables for optimal RAG context."""
    import pandas as pd

    logger.info(f"Converting Excel spreadsheet '{filename}' to Markdown tables...")

    elements = []
    try:
        # Load all sheets
        dict_df = pd.read_excel(io.BytesIO(file_bytes), sheet_name=None)

        for sheet_name, df in dict_df.items():
            if df.empty:
                continue

            # Use pandas markdown conversion for high-quality table representation
            md_table = df.to_markdown(index=False)
            elements.append(
                {"type": "table", "text": f"Sheet: {sheet_name}\n\n{md_table}", "page": 1}
            )
    except Exception as e:
        logger.error(f"Excel parsing failed for {filename}: {e}")
        # Fallback to csv-like if pandas fails
        pass

    return elements


def _parse_pptx(file_bytes: bytes, filename: str) -> List[dict]:
    """Extracts slide text and titles from PowerPoint presentations."""
    from pptx import Presentation

    logger.info(f"Extracting presentation text from '{filename}'...")

    prs = Presentation(io.BytesIO(file_bytes))
    elements = []

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        slide_texts = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Try to identify titles
                if shape == slide.shapes.title:
                    elements.append(
                        {"type": "heading", "text": shape.text.strip(), "page": slide_num}
                    )
                else:
                    slide_texts.append(shape.text.strip())

        if slide_texts:
            elements.append(
                {"type": "paragraph", "text": "\n".join(slide_texts), "page": slide_num}
            )

    return elements


def _parse_text_generic(file_bytes: bytes, filename: str, ext: str) -> List[dict]:
    """Handles all text-based formats with resilient decoding."""
    logger.info(f"Using generic text parser for '{filename}' (Format: {ext})...")

    try:
        # Try UTF-8 first, fallback to Latin-1 for high resilience
        try:
            text = file_bytes.decode("utf-8")
        except UnicodeDecodeError:
            text = file_bytes.decode("latin-1")

        if not text.strip():
            return []

        # Split into elements by paragraph for better indexing
        paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
        return [{"type": "paragraph", "text": p, "page": 1} for p in paragraphs]

    except Exception as e:
        logger.error(f"Generic text parsing failed for {filename}: {e}")
        return []


# ─────────────────────────────────────────────
# Embeddings (nvidia/nv-embedqa-e5-v5)
# ─────────────────────────────────────────────


async def get_embeddings(texts: List[str], input_type: str = "passage") -> List[List[float]]:
    """
    Generate embeddings for a list of texts using NVIDIA NIM embedding model.

    Args:
        texts: List of strings to embed.
        input_type: 'passage' for indexing, 'query' for search queries.

    Returns:
        List of embedding vectors.
    """
    client = _get_nim_client()

    # NIM embedding API uses extra_body for input_type
    response = await client.embeddings.create(
        model=settings.NIM_EMBEDDING_MODEL,
        input=texts,
        extra_body={"input_type": input_type, "truncate": "END"},
    )

    embeddings = [item.embedding for item in response.data]
    logger.info(f"Generated {len(embeddings)} embeddings (input_type={input_type})")
    return embeddings


async def get_query_embedding(query: str) -> List[float]:
    """Convenience wrapper for single query embedding."""
    embeddings = await get_embeddings([query], input_type="query")
    return embeddings[0]


# ─────────────────────────────────────────────
# Reranking (nvidia/nv-rerankqa-mistral-4b-v3)
# ─────────────────────────────────────────────


async def rerank_chunks(query: str, chunks: List[dict], top_n: int = 10) -> List[dict]:
    """
    Rerank a list of chunks based on their relevance to the query using NVIDIA Reranker NIM.

    Args:
        query: The user's search query.
        chunks: List of chunk dictionaries from vector search.
        top_n: Number of top reranked chunks to return.

    Returns:
        The top_n reranked chunks with updated scores.
    """
    if not chunks:
        return []

    url = settings.NIM_RERANK_URL
    headers = {
        "Authorization": f"Bearer {settings.NVIDIA_API_KEY}",
        "Content-Type": "application/json",
        "Accept": "application/json",
    }

    # Prepare passages for NIM
    # The NIM expects a list of dictionaries with 'text'
    passages = [{"text": c["text"]} for c in chunks]

    payload = {
        "model": settings.NIM_RERANK_MODEL,
        "query": {"text": query},
        "passages": passages,
        "truncate": "END",
    }

    async with httpx.AsyncClient() as client:
        try:
            response = await client.post(url, headers=headers, json=payload, timeout=30.0)
            response.raise_for_status()
            data = response.json()

            # The response usually contains 'rankings' with 'index' and 'logit' (score)
            rankings = data.get("rankings", [])
            if not rankings:
                return chunks[:top_n]

            # Extract logits for normalization
            logits = [r["logit"] for r in rankings]
            max_logit = max(logits)
            min_logit = min(logits)
            logit_range = max_logit - min_logit

            reranked_chunks = []
            for rank in rankings[:top_n]:
                idx = rank["index"]
                logit = rank["logit"]

                # Relative Normalization (0.0 to 1.0)
                # If all logits are same, default to 1.0
                if logit_range > 0:
                    score = (logit - min_logit) / logit_range
                    # Soft calibration: Scale to 50%-95% range for the top results to look better in UI
                    score = 0.5 + (score * 0.45)
                else:
                    score = 0.95

                # Update chunk
                chunk = chunks[idx].copy()
                chunk["rerank_score"] = score
                chunk["score"] = score
                reranked_chunks.append(chunk)

            logger.info(
                f"Reranked {len(chunks)} chunks. Raw logits: {logits[:5]}. Normalized scores: {[round(c['score'], 2) for c in reranked_chunks[:5]]}"
            )
            return reranked_chunks

        except Exception as e:
            logger.error(f"NIM Rerank failed: {e}")
            # Fallback: return original chunks if reranking fails
            return chunks[:top_n]


# ─────────────────────────────────────────────
# Chat / RAG Generation (LLM via NIM)
# ─────────────────────────────────────────────


async def generate_rag_answer(
    query: str,
    context_chunks: List[dict],
    system_prompt: Optional[str] = None,
    model_override: Optional[str] = None,
    stream: bool = False,
) -> str:
    """
    Generate a grounded answer using retrieved context chunks and NVIDIA NIM LLM.

    Args:
        query: The user's question.
        context_chunks: List of retrieved chunk dicts (must have 'text', 'filename').
        system_prompt: Optional override for the system prompt.
        model_override: Optional override for the NIM LLM model.
        stream: If True, uses SSE streaming (returns async generator).

    Returns:
        The generated answer string.
    """
    client = _get_nim_client()
    model = model_override or settings.NIM_LLM_MODEL

    # Build RAG context block
    context_parts = []
    for i, chunk in enumerate(context_chunks, 1):
        source = chunk.get("filename", "Unknown Source")
        page = chunk.get("page")
        page_str = f", page {page}" if page else ""
        context_parts.append(f"[Source {i}: {source}{page_str}]\n{chunk['text']}")
    context_str = "\n\n---\n\n".join(context_parts)

    if system_prompt is None:
        system_prompt = (
            "You are an elite AI Research Analyst at a top-tier strategy consulting firm (McKinsey/KPMG style). "
            "Your objective is to provide high-density, data-driven intelligence briefings.\n\n"
            "Executive Guidelines:\n"
            "1. OBJECTIVITY & SYNTHESIS: Maintain a neutral, professional tone. Integrate findings from multiple sources into cohesive paragraphs.\n"
            "2. ATTRIBUTION: Use strict Markdown link notation: [Source N](#source-n) immediately following factual claims.\n"
            "3. ANALYTICS PROTOCOL: When asked for sentiment, provide a structured 'Sentiment Scorecard' (e.g., Sentiment Score: X/10, Primary Emotion: [Emotion], Confidence: Y%).\n"
            "4. VISUALIZATION PROTOCOL: You have access to professional charting engines. If a data visualization would improve interpretation:\n"
            '   - Use ```json:chart ``` code blocks for Data Charts (Bar/Line/Pie). Format: { "type": "bar"|"line"|"pie", "data": [{"name": "Category", "value": 10}], "options": {"title": "Title"} }.\n'
            "   - ALWAYS provide a text-based executive summary of any generated visualization.\n"
            "5. TABULAR SYNTHESIS: When presenting cost estimates, line-item comparisons, or structured thematic areas, ALWAYS use standard Markdown Tables with pipe separators (|) and header delimiters (|---|).\n"
            "6. INTEGRITY: If information is missing, state it explicitly. Do not speculate."
        )

    messages = [
        {"role": "system", "content": system_prompt},
        {
            "role": "user",
            "content": (
                f"I have retrieved the following technical context to help answer my question.\n\n"
                f"### Context:\n{context_str}\n\n"
                f"### Question:\n{query}\n\n"
                f"### Instructions:\n"
                f"Based on the context above, provide a detailed and synthesized response to the question. "
                f"Identify key facts, explain relationships, and ensure all claims are grounded in the 'Sources'.\n\n"
                f"Answer:"
            ),
        },
    ]

    logger.debug(
        f"RAG Prompt Construction (first 500 chars): {str(messages[1]['content'])[:500]}..."
    )

    if stream:
        # Return async generator for SSE
        async def _stream():
            try:
                async with client.chat.completions.with_options(timeout=30.0).stream(
                    model=model,
                    messages=messages,
                    temperature=0.1,
                    max_tokens=2048,
                ) as stream_ctx:
                    async for chunk in stream_ctx:
                        delta = chunk.choices[0].delta.content
                        if delta:
                            yield delta
            except Exception as e:
                logger.error(f"Error streaming RAG answer with model '{model}': {e}")
                yield f"\n\n[Error]: The AI model is currently unavailable (Status: {getattr(e, 'status_code', 'Unknown')}). Please try again in a few moments."

        return _stream()

    try:
        response = await client.chat.completions.create(
            model=model,
            messages=messages,
            temperature=0.1,
            max_tokens=2048,
            timeout=30.0,
        )
        answer = response.choices[0].message.content
        logger.info(f"Generated RAG answer ({len(answer)} chars) for query: '{query[:60]}...'")
        return answer
    except Exception as e:
        logger.error(f"Error generating RAG answer with model '{model}': {e}")
        return f"AI Generation Error: The model '{model}' is currently degraded or unreachable. Please try again later."
